# -*- coding: utf-8 -*-
"""
trust-qna/references 업로드 자료 본문 추출기.
PDF/DOCX/XLSX/PPTX/TXT 본문을 텍스트로 뽑아 references/_extracted/ 아래 .md 로 저장.
build-backdata-index.mjs 가 references 하위 .md/.txt 를 자동 인제스트하므로
_extracted 산출물이 곧 상담 RAG 근거가 된다.
- HWP/구포맷(.doc/.xls)/이미지(.jpg/.png) 미지원 → 스킵(매니페스트 기록)
- 대외비 보호: 원본 절대경로를 본문에 쓰지 않음(파일명만). _extracted 도 git/배포 제외(references/* 통째 ignore).
실행: python scripts/extract-refs.py
"""
import os, sys, datetime

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # trust-qna
SRC = os.path.join(BASE, "references")
OUT = os.path.join(SRC, "_extracted")

PAGE_CAP = 150
CHAR_CAP = 400_000
ROW_CAP = 500
TEXT_EXT = {'.txt', '.md', '.csv'}
SKIP_NAMES = {'readme.md'}

def longp(p):
    p = os.path.abspath(p)
    if os.name == 'nt' and not p.startswith('\\\\?\\'):
        p = '\\\\?\\' + p
    return p

def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    out = []
    n = min(len(doc), PAGE_CAP)
    for i in range(n):
        t = doc[i].get_text()
        if t.strip():
            out.append(f"\n--- p.{i+1} ---\n{t}")
    if len(doc) > PAGE_CAP:
        out.append(f"\n[... {len(doc)-PAGE_CAP} pages 생략 (총 {len(doc)}p) ...]")
    doc.close()
    return "".join(out)

def extract_docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for ti, tbl in enumerate(d.tables):
        parts.append(f"\n[표 {ti+1}]")
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"\n=== 시트: {ws.title} ===")
        r = 0
        for row in ws.iter_rows(values_only=True):
            vals = [("" if v is None else str(v)) for v in row]
            if any(v.strip() for v in vals):
                parts.append("\t".join(vals))
            r += 1
            if r >= ROW_CAP:
                parts.append(f"[... 행 생략 ({ws.max_row}행) ...]")
                break
    wb.close()
    return "\n".join(parts)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for si, slide in enumerate(prs.slides):
        if si >= PAGE_CAP:
            parts.append("[... 슬라이드 생략 ...]")
            break
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs)
                    if line.strip():
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            parts.append(f"\n--- slide {si+1} ---\n" + "\n".join(texts))
    return "\n".join(parts)

def extract_text(path):
    for enc in ('utf-8', 'cp949', 'utf-16'):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(path, encoding='utf-8', errors='replace') as f:
        return f.read()

HANDLERS = {
    '.pdf': extract_pdf, '.docx': extract_docx,
    '.xlsx': extract_xlsx, '.xlsm': extract_xlsx,
    '.pptx': extract_pptx,
}

def main():
    os.makedirs(OUT, exist_ok=True)
    manifest = [f"# trust-qna 추출 매니페스트", "",
                f"> 생성 {datetime.date.today().isoformat()}", "",
                "| 파일 | 상태 | 글자수 |", "|---|---|---:|"]
    done = skipped = failed = unsupported = 0
    for dirpath, dirs, files in os.walk(SRC):
        if os.path.abspath(dirpath).startswith(os.path.abspath(OUT)):
            continue  # 산출물 폴더는 재귀 제외
        for fn in sorted(files):
            if fn.lower() in SKIP_NAMES:
                continue
            ext = os.path.splitext(fn)[1].lower()
            src = os.path.join(dirpath, fn)
            rel = os.path.relpath(src, SRC)
            handler = HANDLERS.get(ext)
            if handler is None and ext not in TEXT_EXT:
                unsupported += 1
                manifest.append(f"| {rel} | ⏭️ 미지원({ext}) | - |")
                continue
            try:
                text = (handler(src) if handler else extract_text(src)) or ""
                text = text.strip()
            except Exception as e:
                failed += 1
                manifest.append(f"| {rel} | ❌ {type(e).__name__} | - |")
                continue
            if not text:
                skipped += 1
                manifest.append(f"| {rel} | ⚠️ 빈 텍스트(스캔본?) | 0 |")
                continue
            if len(text) > CHAR_CAP:
                text = text[:CHAR_CAP] + f"\n\n[... {CHAR_CAP} 초과 절단 ...]"
            # 출력: 원본경로 없이 파일명만(누출 방지). 평탄화하되 경로 충돌 방지 위해 rel 기반 안전 파일명.
            safe = rel.replace('\\', '__').replace('/', '__')
            out_path = os.path.join(OUT, safe) + ".md"
            try:
                with open(longp(out_path), 'w', encoding='utf-8') as f:
                    f.write(f"# {fn}\n\n{text}")
            except Exception as e:
                failed += 1
                manifest.append(f"| {rel} | ❌ write:{type(e).__name__} | - |")
                continue
            done += 1
            manifest.append(f"| {rel} | ✅ | {len(text):,} |")
    manifest.insert(4, f"**완료 {done} · 빈본문 {skipped} · 미지원 {unsupported} · 실패 {failed}**\n")
    with open(os.path.join(OUT, "_manifest.md"), 'w', encoding='utf-8') as f:
        f.write("\n".join(manifest))
    print(f"done={done} empty={skipped} unsupported={unsupported} fail={failed}", flush=True)

if __name__ == '__main__':
    main()
